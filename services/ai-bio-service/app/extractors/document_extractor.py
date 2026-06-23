import re
from dataclasses import dataclass
from io import BytesIO

import markdown as markdown_lib
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
from pypdf import PdfReader

from app.core.error_codes import ErrorCode
from app.sources.source_type import SourceType


class DocumentExtractionError(Exception):
    def __init__(
        self,
        *,
        code: ErrorCode | str,
        message: str,
        retryable: bool = False,
    ) -> None:
        self.code = str(code)
        self.message = message
        self.retryable = retryable
        super().__init__(message)


@dataclass(frozen=True)
class ExtractionResult:
    extracted_text: str
    cleaned_text: str
    extracted_char_count: int
    cleaned_char_count: int
    parser_name: str
    parser_version: str
    warnings: list[dict]


class DocumentExtractor:
    parser_version = "phase1-v1"

    def extract(
        self,
        *,
        source_type: SourceType,
        content: bytes,
    ) -> ExtractionResult:
        if source_type == SourceType.PDF:
            return self._extract_pdf(content)

        if source_type == SourceType.TEXT:
            return self._extract_text(content)

        if source_type == SourceType.MARKDOWN:
            return self._extract_markdown(content)

        if source_type == SourceType.DOCX:
            return self._extract_docx(content)

        if source_type == SourceType.PPTX:
            return self._extract_pptx(content)

        raise DocumentExtractionError(
            code=ErrorCode.UNSUPPORTED_SOURCE_TYPE,
            message="Source type is not supported by the Phase 1 extractor.",
            retryable=False,
        )

    def _extract_pdf(self, content: bytes) -> ExtractionResult:
        try:
            reader = PdfReader(BytesIO(content))
        except Exception as exc:
            raise DocumentExtractionError(
                code=ErrorCode.INVALID_SOURCE_TYPE,
                message="PDF structure is invalid.",
                retryable=False,
            ) from exc

        if reader.is_encrypted:
            raise DocumentExtractionError(
                code=ErrorCode.DOCUMENT_PASSWORD_PROTECTED,
                message="Password-protected PDF is not supported.",
                retryable=False,
            )

        warnings: list[dict] = []
        pages: list[str] = []

        for index, page in enumerate(reader.pages):
            try:
                page_text = page.extract_text() or ""
            except Exception:
                page_text = ""
                warnings.append(
                    {
                        "page": index + 1,
                        "message": "Page text extraction failed.",
                    }
                )

            if page_text.strip():
                pages.append(page_text)

        extracted_text = "\n\n".join(pages)

        return self._build_result(
            extracted_text=extracted_text,
            parser_name="pypdf",
            warnings=warnings,
        )

    def _extract_text(self, content: bytes) -> ExtractionResult:
        try:
            extracted_text = content.decode("utf-8")
        except UnicodeDecodeError as exc:
            raise DocumentExtractionError(
                code=ErrorCode.INVALID_SOURCE_TYPE,
                message="Text source must be valid UTF-8.",
                retryable=False,
            ) from exc

        return self._build_result(
            extracted_text=extracted_text,
            parser_name="utf8-text",
            warnings=[],
        )

    def _extract_markdown(self, content: bytes) -> ExtractionResult:
        try:
            markdown_text = content.decode("utf-8")
        except UnicodeDecodeError as exc:
            raise DocumentExtractionError(
                code=ErrorCode.INVALID_SOURCE_TYPE,
                message="Markdown source must be valid UTF-8.",
                retryable=False,
            ) from exc

        html = markdown_lib.markdown(markdown_text)
        extracted_text = BeautifulSoup(html, "html.parser").get_text("\n")

        return self._build_result(
            extracted_text=extracted_text,
            parser_name="markdown-beautifulsoup",
            warnings=[],
        )

    def _extract_docx(self, content: bytes) -> ExtractionResult:
        try:
            document = Document(BytesIO(content))
        except Exception as exc:
            raise DocumentExtractionError(
                code=ErrorCode.INVALID_SOURCE_TYPE,
                message="DOCX structure is invalid.",
                retryable=False,
            ) from exc

        parts: list[str] = []

        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if text:
                parts.append(text)

        for table in document.tables:
            for row in table.rows:
                row_text = " | ".join(
                    cell.text.strip()
                    for cell in row.cells
                    if cell.text.strip()
                )
                if row_text:
                    parts.append(row_text)

        return self._build_result(
            extracted_text="\n".join(parts),
            parser_name="python-docx",
            warnings=[],
        )

    def _extract_pptx(self, content: bytes) -> ExtractionResult:
        try:
            presentation = Presentation(BytesIO(content))
        except Exception as exc:
            raise DocumentExtractionError(
                code=ErrorCode.INVALID_SOURCE_TYPE,
                message="PPTX structure is invalid.",
                retryable=False,
            ) from exc

        parts: list[str] = []

        for slide_index, slide in enumerate(presentation.slides):
            slide_parts: list[str] = []

            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    slide_parts.append(text.strip())

            if slide_parts:
                parts.append(f"Slide {slide_index + 1}\n" + "\n".join(slide_parts))

        return self._build_result(
            extracted_text="\n\n".join(parts),
            parser_name="python-pptx",
            warnings=[],
        )

    def _build_result(
        self,
        *,
        extracted_text: str,
        parser_name: str,
        warnings: list[dict],
    ) -> ExtractionResult:
        cleaned_text = self._basic_clean(extracted_text)

        if not cleaned_text:
            raise DocumentExtractionError(
                code=ErrorCode.NO_USABLE_SOURCE_CONTENT,
                message="No usable text content was found in the source.",
                retryable=False,
            )

        return ExtractionResult(
            extracted_text=extracted_text,
            cleaned_text=cleaned_text,
            extracted_char_count=len(extracted_text),
            cleaned_char_count=len(cleaned_text),
            parser_name=parser_name,
            parser_version=self.parser_version,
            warnings=warnings,
        )

    def _basic_clean(self, text: str) -> str:
        text = text.replace("\x00", " ")
        text = re.sub(r"[ \t]+", " ", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()


document_extractor = DocumentExtractor()